# backend/doc_loader.py
import re
from pathlib import Path
from typing import List, Dict, Optional
from langchain_core.documents import Document
try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
import PyPDF2

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    docx = None
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    PPTX_AVAILABLE = False

from .config import Config

class UniversalDocLoader:
    @staticmethod
    def load_document(file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        if ext == ".pdf":
            return UniversalDocLoader._load_pdf(file_path)
        if ext == ".txt":
            return UniversalDocLoader._load_txt(file_path)
        if ext == ".docx":
            return UniversalDocLoader._load_docx(file_path)
        if ext == ".pptx":
            return UniversalDocLoader._load_pptx(file_path)
        raise ValueError(f"Unsupported format: {ext}")

    @staticmethod
    def _load_pdf(path: str) -> str:
        text = []
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text() or ""
                text.append(page_text)
        return "\n".join(text)

    @staticmethod
    def _load_txt(path: str) -> str:
        return Path(path).read_text(encoding="utf-8", errors="ignore")

    @staticmethod
    def _load_docx(path: str) -> str:
        if not DOCX_AVAILABLE or docx is None:
            raise ImportError("Install python-docx")
        doc = docx.Document(path)
        # Filter out empty paragraphs and join with newlines
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)

    @staticmethod
    def _load_pptx(path: str) -> str:
        if not PPTX_AVAILABLE or Presentation is None:
            raise ImportError("Install python-pptx")
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape_text = getattr(shape, "text", "").strip()
                    if shape_text:
                        text.append(shape_text)
        return "\n".join(text)

    @staticmethod
    def clean_text(text: str) -> str:
        t = re.sub(r"\t", " ", text)
        t = re.sub(r"\n\s*\n", "\n", t)
        t = re.sub(r"(\w)\n(\w)", r"\1 \2", t)
        t = re.sub(r" +", " ", t)
        return t.strip()

    @staticmethod
    def chunk_text(text: str, metadata: Optional[Dict[str, str]] = None) -> List[Document]:
        if metadata is None:
            metadata = {}
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=Config.CHUNK_SIZE,
            chunk_overlap=Config.CHUNK_OVERLAP
        )
        chunks = splitter.split_text(text)
        documents = []
        for i, chunk in enumerate(chunks, 1):
            chunk_metadata = metadata.copy()
            chunk_metadata["chunk_id"] = str(i)
            documents.append(Document(page_content=chunk, metadata=chunk_metadata))
        return documents

    @staticmethod
    def chunk_pdf_by_pages(file_path: str, metadata: Optional[Dict[str, str]] = None) -> List[Document]:
        """Load PDF and chunk by pages for better citation accuracy"""
        if metadata is None:
            metadata = {}

        documents = []
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(reader.pages, 1):
                page_text = page.extract_text() or ""
                if page_text.strip():  # Only process non-empty pages
                    page_text = UniversalDocLoader.clean_text(page_text)

                    # If page text is too long, chunk it further
                    if len(page_text) > Config.CHUNK_SIZE:
                        splitter = RecursiveCharacterTextSplitter(
                            chunk_size=Config.CHUNK_SIZE,
                            chunk_overlap=Config.CHUNK_OVERLAP
                        )
                        chunks = splitter.split_text(page_text)
                        for chunk_num, chunk in enumerate(chunks, 1):
                            chunk_metadata = metadata.copy()
                            chunk_metadata["page_number"] = str(page_num)
                            chunk_metadata["chunk_id"] = f"{page_num}.{chunk_num}"
                            documents.append(Document(page_content=chunk, metadata=chunk_metadata))
                    else:
                        # Page fits in single chunk
                        chunk_metadata = metadata.copy()
                        chunk_metadata["page_number"] = str(page_num)
                        chunk_metadata["chunk_id"] = str(page_num)
                        documents.append(Document(page_content=page_text, metadata=chunk_metadata))

        return documents
